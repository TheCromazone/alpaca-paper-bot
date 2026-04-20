"""Targeted text updates for the Claude Design deck.

Edits specific runs in place so we keep the design's exact fonts, colors, and
positioning. We only touch text — no layout changes. Re-run any time the bot's
reality drifts from the deck's claims.
"""
import sys
from pathlib import Path
from pptx import Presentation

sys.stdout.reconfigure(encoding="utf-8")

# NB: The Claude Design deck uses non-breaking hyphens (U+2011) in many places.
# Target strings here must match the deck's exact Unicode.

REPLACEMENTS = {
    # ---- Slide 5: At a glance ----
    "83": "102",
    "75 equities + 8 bond ETFs": "94 equities + 8 bond ETFs",

    # ---- Slide 7: News sentiment signal ----
    (
        "Every 5 min we pull RSS feeds from six financial outlets (CNBC, Reuters,"
        " Yahoo Finance, CNN Business, MarketWatch, Seeking Alpha). Each headline is"
        " scored with VADER (a lexicon\u2011based sentiment model). We match"
        " headlines to tickers, take a 7\u2011day rolling average, z\u2011score it."
        " Optionally re\u2011score with FinBERT, a finance\u2011tuned transformer,"
        " for headlines that mention a universe ticker."
    ): (
        "Every 5 min we pull RSS feeds from five financial outlets (CNBC, Yahoo"
        " Finance, CNN Business, MarketWatch, Seeking Alpha). Each headline is scored"
        " with VADER. For headlines that tag a universe ticker we also scrape the"
        " full article body with trafilatura and re\u2011score per\u2011sentence,"
        " then take a 7\u2011day rolling z\u2011score. Optional FinBERT re\u2011score"
        " layered on top for the short\u2011list."
    ),
    "CNBC, Reuters, Yahoo Finance, CNN Business, MarketWatch, Seeking Alpha":
        "CNBC, Yahoo Finance, CNN Business, MarketWatch, Seeking Alpha",
    "VADER polarity \u2192 7d rolling z\u2011score (FinBERT optional)":
        "headline + scraped article body \u2192 per\u2011sentence VADER \u2192 7d z\u2011score",

    # ---- Slide 19: Limits (drop the "headlines only" limit; article scraping is live) ----
    "Headlines only, no article bodies":
        "No paywalled wires",
    (
        "We only read RSS, not the articles behind them. A careful headline"
        " can miss crucial context in paragraph four."
    ): (
        "We fetch full article bodies for every headline that tags a universe"
        " ticker (CNBC, Yahoo, MarketWatch, CNN Business, Seeking Alpha) via"
        " trafilatura. Paywalled outlets like Bloomberg and WSJ remain out of reach."
    ),
    "\u2192 Unavoidable without paid licenses; VADER on headlines still has signal.":
        "\u2192 Body\u2011level VADER is a meaningful upgrade over headlines alone.",

    # ---- Slide 21: Roadmap (update ticker-cap line) ----
    (
        "Lift the 83\u2011ticker cap; admit S&P 1500 + liquid ETFs, gated by a"
        " liquidity filter so IEX sparseness stays tolerable."
    ): (
        "Lift the current ~100\u2011ticker cap; admit S&P 1500 + liquid ETFs,"
        " gated by a liquidity filter so IEX sparseness stays tolerable."
    ),
}


def replace_in_paragraph(para, old: str, new: str) -> bool:
    """Replace `old` with `new` inside a paragraph. Returns True if it fired.

    A paragraph's text is split across `runs` (one per formatting change). The
    target string may span multiple runs, so we concatenate, substitute, and
    then write back — placing the full replacement into the first run, and
    clearing the others to preserve baseline formatting of the first run.
    """
    joined = "".join(r.text for r in para.runs)
    if old not in joined:
        return False
    new_joined = joined.replace(old, new)
    if not para.runs:
        return False
    para.runs[0].text = new_joined
    for r in para.runs[1:]:
        r.text = ""
    return True


def main():
    path = Path(__file__).resolve().parent.parent / "Alpaca Bot Deck.pptx"
    p = Presentation(str(path))
    hits = {k: 0 for k in REPLACEMENTS}

    for slide in p.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for old, new in REPLACEMENTS.items():
                    if replace_in_paragraph(para, old, new):
                        hits[old] += 1

    for old, n in hits.items():
        flag = "\u2713" if n else "\u2717"
        print(f"  [{flag}] {n}x  {old[:70]!r}")

    p.save(str(path))
    print(f"Wrote {path}")


if __name__ == "__main__":
    main()
