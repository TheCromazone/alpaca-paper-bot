import sys
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
p = Presentation("deck/Alpaca Bot Deck.pptx")
for i, slide in enumerate(p.slides, 1):
    if i not in (5, 7, 19, 21):
        continue
    print(f"=== Slide {i} ===")
    for sh in slide.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        print(f"  {run.text!r}")
